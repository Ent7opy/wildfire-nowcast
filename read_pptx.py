from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE_TYPE

prs = Presentation(r'C:\Users\vanyo\Desktop\Wildfire Nowcast Mid Demo.pptx')

print(f"Total slides: {len(prs.slides)}")
print(f"Slide width: {prs.slide_width.inches:.2f} inches")
print(f"Slide height: {prs.slide_height.inches:.2f} inches")
print("-" * 60)

for i, slide in enumerate(prs.slides):
    print(f"\n=== SLIDE {i+1} ===")
    layout_name = slide.slide_layout.name if slide.slide_layout else "Unknown"
    print(f"Layout: {layout_name}")

    for shape in slide.shapes:
        if shape.has_text_frame:
            text = shape.text_frame.text.strip()
            if text:
                shape_type = "Title" if shape == slide.shapes.title else "Text"
                print(f"  [{shape_type}] {text[:100]}{'...' if len(text) > 100 else ''}")
        elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            print(f"  [Image] {shape.width.inches:.1f}x{shape.height.inches:.1f} inches")
        elif shape.shape_type == MSO_SHAPE_TYPE.TABLE:
            print(f"  [Table] {shape.table.rows} rows x {shape.table.columns} cols")
