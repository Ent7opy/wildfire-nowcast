from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# Load the presentation
prs = Presentation(r'C:\Users\vanyo\Desktop\Wildfire Nowcast Mid Demo.pptx')

# Get the layouts we can use
layout_subtitle = prs.slides[1].slide_layout  # 02_Subtitle_Slide_02
layout_title_text = prs.slides[2].slide_layout  # 07_Title&text

# Helper function to add text box with styling
def add_text_box(slide, left, top, width, height, text, font_size=14, bold=False, color=None, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.alignment = align
    if color:
        p.font.color.rgb = RGBColor(*color)
    return txBox

def add_bullet_text(slide, left, top, width, height, items, font_size=12, title=None, title_size=16):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True

    first = True
    if title:
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(title_size)
        p.font.bold = True
        first = False

    for item in items:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(font_size)
        p.level = 0
    return txBox

def add_box_shape(slide, left, top, width, height, text, fill_color, font_size=10, text_color=(255,255,255)):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(*fill_color)
    shape.line.fill.background()

    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = True
    p.font.color.rgb = RGBColor(*text_color)
    p.alignment = PP_ALIGN.CENTER

    # Center vertically
    shape.text_frame.auto_size = None

    return shape

# Define colors (modern palette)
COLOR_UI = (66, 133, 244)       # Blue - UI
COLOR_API = (52, 168, 83)       # Green - API
COLOR_DB = (251, 188, 4)        # Yellow - Database
COLOR_CACHE = (234, 67, 53)     # Red - Redis
COLOR_TILES = (103, 58, 183)    # Purple - Tile servers
COLOR_WORKER = (255, 109, 0)    # Orange - Worker
COLOR_ML = (0, 150, 136)        # Teal - ML

# ============================================
# SLIDE 2: System Architecture Overview
# ============================================
slide2 = prs.slides[1]

# Add title
add_text_box(slide2, 0.5, 0.3, 9, 0.6, "System Architecture", font_size=28, bold=True, align=PP_ALIGN.CENTER)
add_text_box(slide2, 0.5, 0.75, 9, 0.4, "7 Integrated Services", font_size=16, align=PP_ALIGN.CENTER, color=(100, 100, 100))

# Row 1: UI Layer
add_box_shape(slide2, 3.5, 1.3, 3, 0.6, "UI (Streamlit) :8501\nMap & Filters", COLOR_UI, 9)

# Row 2: API Layer
add_box_shape(slide2, 3.5, 2.2, 3, 0.6, "API (FastAPI) :8000\nREST Endpoints", COLOR_API, 9)

# Row 3: Data Layer (3 boxes)
add_box_shape(slide2, 0.8, 3.1, 2.2, 0.7, "PostgreSQL\n+ PostGIS\n:5432", COLOR_DB, 8, (50, 50, 50))
add_box_shape(slide2, 3.5, 3.1, 2.2, 0.7, "Redis\nCache + Queue\n:6379", COLOR_CACHE, 8)
add_box_shape(slide2, 6.2, 3.1, 2.2, 0.7, "TiTiler\nCOG Tiles\n:8080", COLOR_TILES, 8)

# Row 4: Workers and Vector Tiles
add_box_shape(slide2, 2.0, 4.1, 2.5, 0.6, "Background Worker\nRQ Jobs", COLOR_WORKER, 9)
add_box_shape(slide2, 5.5, 4.1, 2.5, 0.6, "pg_tileserv\nVector Tiles :7800", COLOR_TILES, 9)

# Add connection lines (simple text arrows)
add_text_box(slide2, 4.7, 1.9, 0.6, 0.3, "↓", font_size=16, align=PP_ALIGN.CENTER)
add_text_box(slide2, 2.5, 2.7, 0.6, 0.3, "↓", font_size=16, align=PP_ALIGN.CENTER)
add_text_box(slide2, 4.7, 2.7, 0.6, 0.3, "↓", font_size=16, align=PP_ALIGN.CENTER)
add_text_box(slide2, 7.0, 2.7, 0.6, 0.3, "↓", font_size=16, align=PP_ALIGN.CENTER)
add_text_box(slide2, 3.0, 3.7, 0.6, 0.3, "↓", font_size=16, align=PP_ALIGN.CENTER)

# Add legend on the right
add_text_box(slide2, 8.7, 1.3, 1.3, 0.3, "Data Flow:", font_size=10, bold=True)
add_text_box(slide2, 8.7, 1.55, 1.5, 1.8,
    "1. Ingest pipelines\n   fetch external data\n\n"
    "2. Worker processes\n   & stores in DB\n\n"
    "3. API serves data\n   to UI\n\n"
    "4. TiTiler serves\n   raster tiles", font_size=8)

# ============================================
# SLIDE 3: Data Pipeline
# ============================================
slide3 = prs.slides[2]

add_text_box(slide3, 0.5, 0.3, 9, 0.6, "Data Pipeline & Sources", font_size=28, bold=True, align=PP_ALIGN.CENTER)

# Pipeline flow boxes
pipeline_y = 1.2
box_width = 1.3
box_height = 0.5
gap = 0.15
start_x = 0.5

steps = [
    ("INGEST", COLOR_WORKER),
    ("DENOISE", COLOR_ML),
    ("FEATURES", COLOR_API),
    ("FORECAST", COLOR_ML),
    ("CALIBRATE", COLOR_ML),
    ("SERVE", COLOR_API),
]

for i, (label, color) in enumerate(steps):
    x = start_x + i * (box_width + gap)
    add_box_shape(slide3, x, pipeline_y, box_width, box_height, label, color, 9)
    if i < len(steps) - 1:
        add_text_box(slide3, x + box_width, pipeline_y + 0.1, 0.2, 0.3, "→", font_size=14, align=PP_ALIGN.CENTER)

# Data sources section
add_text_box(slide3, 0.5, 2.0, 4, 0.4, "External Data Sources", font_size=14, bold=True)

sources = [
    "NASA FIRMS - Fire detections (VIIRS, MODIS)",
    "NOAA GFS - Weather forecasts (wind, temp, humidity)",
    "Copernicus DEM - Terrain (30m elevation, slope)",
    "ERA5 - Reanalysis for bias correction (optional)"
]
add_bullet_text(slide3, 0.5, 2.3, 4.5, 1.5, sources, font_size=11)

# Storage section
add_text_box(slide3, 5.2, 2.0, 4.5, 0.4, "What Gets Stored", font_size=14, bold=True)

storage = [
    "fire_detections - with denoised scores",
    "weather_runs - GFS forecast metadata",
    "spread_forecasts - probability rasters (COG)",
    "aois - user-defined areas of interest"
]
add_bullet_text(slide3, 5.2, 2.3, 4.5, 1.5, storage, font_size=11)

# API endpoints
add_text_box(slide3, 0.5, 4.0, 9, 0.4, "Key API Endpoints", font_size=14, bold=True)
endpoints_text = "/fires - Query detections  •  /forecast - Spread predictions  •  /risk - Fire risk maps  •  /tiles - Raster tiles  •  /exports - PNG/CSV/GeoJSON"
add_text_box(slide3, 0.5, 4.35, 9, 0.4, endpoints_text, font_size=11, align=PP_ALIGN.CENTER)

# ============================================
# SLIDE 4: AI/ML Current Implementation (NEW)
# ============================================
slide4 = prs.slides.add_slide(layout_title_text)

add_text_box(slide4, 0.5, 0.3, 9, 0.6, "AI/ML - Current Implementation", font_size=28, bold=True, align=PP_ALIGN.CENTER)

# Hotspot Denoiser
add_box_shape(slide4, 0.5, 1.1, 4.3, 0.45, "Hotspot Denoiser", COLOR_ML, 12)
denoiser_items = [
    "Filters false positives (noise, industrial heat)",
    "Model: RandomForest / XGBoost classifier",
    "Features: brightness, FRP, spatial context",
    "Output: denoised_score [0-1] + is_noise flag",
    "Status: In Production"
]
add_bullet_text(slide4, 0.5, 1.6, 4.3, 1.4, denoiser_items, font_size=10)

# Spread Forecasting
add_box_shape(slide4, 5.2, 1.1, 4.3, 0.45, "Spread Forecasting", COLOR_WORKER, 12)
spread_items = [
    "Predicts fire spread at T+24/48/72h",
    "v0 (Default): Heuristic wind+terrain model",
    "v1 (Experimental): HistGradientBoosting",
    "Output: Probability rasters per horizon",
    "Status: Both Implemented"
]
add_bullet_text(slide4, 5.2, 1.6, 4.3, 1.4, spread_items, font_size=10)

# Probability Calibration
add_box_shape(slide4, 0.5, 3.2, 4.3, 0.45, "Probability Calibration", COLOR_API, 12)
calib_items = [
    "Maps raw scores to calibrated probabilities",
    "Ensures 30% prediction = 30% observed",
    "Metrics: Brier score, ECE, reliability",
    "Status: Framework Ready"
]
add_bullet_text(slide4, 0.5, 3.7, 4.3, 1.2, calib_items, font_size=10)

# Weather Bias Correction
add_box_shape(slide4, 5.2, 3.2, 4.3, 0.45, "Weather Bias Correction", COLOR_UI, 12)
weather_items = [
    "Reduces GFS forecast systematic errors",
    "Compares to ERA5 reanalysis (truth)",
    "Per-variable affine correction",
    "Status: Training Pipeline Ready"
]
add_bullet_text(slide4, 5.2, 3.7, 4.3, 1.2, weather_items, font_size=10)

# ============================================
# SLIDE 5: AI/ML Future / Planned (NEW)
# ============================================
slide5 = prs.slides.add_slide(layout_subtitle)

add_text_box(slide5, 0.5, 0.3, 9, 0.6, "AI/ML - Future Roadmap", font_size=28, bold=True, align=PP_ALIGN.CENTER)

# Planned features with status indicators
add_text_box(slide5, 0.5, 1.1, 9, 0.4, "Planned Enhancements", font_size=18, bold=True)

# LLM Summaries
add_box_shape(slide5, 0.5, 1.55, 0.35, 0.35, "1", (100, 100, 100), 12)
add_text_box(slide5, 1.0, 1.55, 4, 0.35, "LLM-Based AOI Summaries", font_size=13, bold=True)
add_text_box(slide5, 1.0, 1.85, 4, 0.5, "Natural language explanations of forecasts with confidence levels, key drivers, and uncertainty caveats", font_size=10)

# Fire Risk Index
add_box_shape(slide5, 0.5, 2.45, 0.35, 0.35, "2", (100, 100, 100), 12)
add_text_box(slide5, 1.0, 2.45, 4, 0.35, "Fire Risk Index", font_size=13, bold=True)
add_text_box(slide5, 1.0, 2.75, 4, 0.5, "Predict new-fire probability using vegetation, weather, terrain, and recent precipitation", font_size=10)

# Multi-Model Ensemble
add_box_shape(slide5, 5.2, 1.55, 0.35, 0.35, "3", (100, 100, 100), 12)
add_text_box(slide5, 5.7, 1.55, 4, 0.35, "Multi-Model Ensemble", font_size=13, bold=True)
add_text_box(slide5, 5.7, 1.85, 4, 0.5, "Combine heuristic v0 + learned v1 + future models for improved accuracy", font_size=10)

# Regional Tuning
add_box_shape(slide5, 5.2, 2.45, 0.35, 0.35, "4", (100, 100, 100), 12)
add_text_box(slide5, 5.7, 2.45, 4, 0.35, "Regional Model Tuning", font_size=13, bold=True)
add_text_box(slide5, 5.7, 2.75, 4, 0.5, "Geography-specific models trained on regional fire behavior patterns", font_size=10)

# Key principles box
add_text_box(slide5, 0.5, 3.5, 9, 0.4, "Design Principles", font_size=18, bold=True)

principles_box = slide5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(3.9), Inches(9), Inches(1.3))
principles_box.fill.solid()
principles_box.fill.fore_color.rgb = RGBColor(245, 245, 245)
principles_box.line.color.rgb = RGBColor(200, 200, 200)

add_text_box(slide5, 0.7, 4.0, 8.6, 1.0,
    "Transparency: Explicit about model limitations & uncertainty\n"
    "Open Data: Uses only public sources (NASA, NOAA, Copernicus)\n"
    "Modular ML: Model contract enables easy swapping of implementations",
    font_size=11)

# ============================================
# SLIDE 6: Tech Stack Summary (NEW)
# ============================================
slide6 = prs.slides.add_slide(layout_title_text)

add_text_box(slide6, 0.5, 0.3, 9, 0.6, "Technology Stack", font_size=28, bold=True, align=PP_ALIGN.CENTER)

# Backend column
add_text_box(slide6, 0.5, 1.0, 3, 0.4, "Backend & API", font_size=14, bold=True, color=COLOR_API)
backend = ["FastAPI + Uvicorn", "SQLAlchemy + Alembic", "PostgreSQL + PostGIS", "Redis + RQ Workers"]
add_bullet_text(slide6, 0.5, 1.35, 3, 1.2, backend, font_size=11)

# ML column
add_text_box(slide6, 3.5, 1.0, 3, 0.4, "ML & Data", font_size=14, bold=True, color=COLOR_ML)
ml = ["scikit-learn + XGBoost", "xarray + rasterio", "NumPy + Pandas", "joblib (serialization)"]
add_bullet_text(slide6, 3.5, 1.35, 3, 1.2, ml, font_size=11)

# Geospatial column
add_text_box(slide6, 6.5, 1.0, 3, 0.4, "Geospatial", font_size=14, bold=True, color=COLOR_TILES)
geo = ["TiTiler (COG tiles)", "pg_tileserv (vectors)", "GDAL + Shapely", "rio-cogeo + cfgrib"]
add_bullet_text(slide6, 6.5, 1.35, 3, 1.2, geo, font_size=11)

# Frontend row
add_text_box(slide6, 0.5, 2.8, 3, 0.4, "Frontend", font_size=14, bold=True, color=COLOR_UI)
frontend = ["Streamlit", "Pydeck (maps)", "Folium"]
add_bullet_text(slide6, 0.5, 3.15, 3, 1.0, frontend, font_size=11)

# DevOps row
add_text_box(slide6, 3.5, 2.8, 3, 0.4, "DevOps", font_size=14, bold=True, color=COLOR_WORKER)
devops = ["Docker Compose", "uv (pkg manager)", "Ruff (linter)"]
add_bullet_text(slide6, 3.5, 3.15, 3, 1.0, devops, font_size=11)

# Single command box
cmd_box = slide6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(4.3), Inches(9), Inches(0.7))
cmd_box.fill.solid()
cmd_box.fill.fore_color.rgb = RGBColor(40, 44, 52)
cmd_box.line.fill.background()

add_text_box(slide6, 0.7, 4.4, 8.6, 0.5, "docker compose up --build   # Starts entire stack", font_size=14, color=(152, 195, 121))

# Save the presentation
output_path = r'C:\Users\vanyo\Desktop\Wildfire Nowcast Mid Demo.pptx'
prs.save(output_path)
print(f"Presentation saved to: {output_path}")
print(f"Total slides: {len(prs.slides)}")
